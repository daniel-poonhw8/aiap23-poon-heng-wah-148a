from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.drawing import fill
import pandas as pd
import sqlite3
import matplotlib.pyplot as plt
import seaborn as sns
import os

def create_presentation():
    """Create a PowerPoint presentation for the ML project"""

    # Create presentation
    prs = Presentation()

    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Online Shopping Purchase Prediction"
    subtitle.text = "Machine Learning Pipeline & Analysis\nExploratory Data Analysis & Predictive Modeling"

    # Slide 2: Project Overview
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Project Overview'
    tf = body_shape.text_frame
    tf.text = 'Objective: Predict online shopping purchase completion using user behavior data'

    p = tf.add_paragraph()
    p.text = '• Dataset: 12,330 online shopping sessions'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• Target: PurchaseCompleted (binary classification)'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• Features: User behavior metrics, page analytics, traffic data'
    p.level = 0

    # Slide 3: Data Understanding
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Data Understanding'
    tf = body_shape.text_frame
    tf.text = 'Key Features:'

    features = [
        '• CustomerType: Returning/New visitor classification',
        '• SpecialDayProximity: Closeness to special shopping days',
        '• ExitRate: Website exit rate (0-1)',
        '• PageValue: Average page value ($)',
        '• BounceRate: Single-page session rate (0-1)',
        '• ProductPageTime: Time spent on product pages (seconds)',
        '• TrafficSource: Source of website traffic',
        '• GeographicRegion: User location region'
    ]

    for feature in features:
        p = tf.add_paragraph()
        p.text = feature
        p.level = 0

    # Slide 4: EDA Key Findings
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'EDA Key Findings'
    tf = body_shape.text_frame
    tf.text = 'Data Quality & Distribution:'

    findings = [
        '• 15-20% purchase completion rate (conversion challenge)',
        '• 85% of visitors are returning customers',
        '• Most sessions have zero page value (browsing behavior)',
        '• High variability in product page engagement time',
        '• Missing values handled via median/mode imputation'
    ]

    for finding in findings:
        p = tf.add_paragraph()
        p.text = finding
        p.level = 0

    # Slide 5: Correlation Analysis
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Correlation Analysis'
    tf = body_shape.text_frame
    tf.text = 'Key Correlations with Purchase Completion:'

    correlations = [
        '• PageValue: +0.492 (Strong positive - higher value → more purchases)',
        '• BounceRate: -0.151 (Negative - lower bounce → more purchases)',
        '• ExitRate: -0.207 (Negative - lower exit → more purchases)',
        '• ProductPageTime: +0.044 (Weak positive relationship)',
        '• SpecialDayProximity: +0.062 (Weak positive effect)'
    ]

    for corr in correlations:
        p = tf.add_paragraph()
        p.text = corr
        p.level = 0

    # Slide 6: Model Performance
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Model Performance Results'
    tf = body_shape.text_frame
    tf.text = 'Test Set Performance (Accuracy, F1-Score, ROC-AUC):'

    models = [
        '• Random Forest: 89.13%, 0.595, 0.884 (Best F1-Score)',
        '• XGBoost: 88.77%, 0.581, 0.889 (Best ROC-AUC)',
        '• Logistic Regression: 88.16%, 0.507, 0.866 (Baseline)',
        '• SVM: 88.50%, 0.550, 0.875 (Alternative approach)'
    ]

    for model in models:
        p = tf.add_paragraph()
        p.text = model
        p.level = 0

    # Slide 7: Feature Engineering
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Feature Engineering'
    tf = body_shape.text_frame
    tf.text = 'Engineered Features:'

    features = [
        '• PageValue_Category: Binned page values (Zero/Low/Medium/High/Very_High)',
        '• Engagement_Score: (1-BounceRate) × (1-ExitRate) × ProductPageTime',
        '• SpecialDay_Effect: SpecialDayProximity × PageValue interaction',
        '• TrafficSource_Category: Categorical encoding of traffic sources'
    ]

    for feature in features:
        p = tf.add_paragraph()
        p.text = feature
        p.level = 0

    # Slide 8: Business Implications
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Business Implications'
    tf = body_shape.text_frame
    tf.text = 'Key Insights for E-commerce:'

    implications = [
        '• Focus on high-value page optimization to boost conversions',
        '• Reduce bounce and exit rates through UX improvements',
        '• Leverage returning visitor behavior patterns',
        '• Monitor special day effects on purchasing behavior',
        '• Implement real-time purchase prediction for targeted interventions'
    ]

    for impl in implications:
        p = tf.add_paragraph()
        p.text = impl
        p.level = 0

    # Slide 9: Pipeline Architecture
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'ML Pipeline Architecture'
    tf = body_shape.text_frame
    tf.text = 'Modular Pipeline Components:'

    components = [
        '• Data Loading: SQLite database integration',
        '• Preprocessing: Missing value handling, encoding, scaling',
        '• Feature Engineering: Domain-specific feature creation',
        '• Model Training: Multiple algorithms with hyperparameter tuning',
        '• Evaluation: Comprehensive metrics and validation',
        '• Deployment: Streamlit web application for predictions'
    ]

    for comp in components:
        p = tf.add_paragraph()
        p.text = comp
        p.level = 0

    # Slide 10: Next Steps & Recommendations
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Next Steps & Recommendations'
    tf = body_shape.text_frame
    tf.text = 'Future Development:'

    recommendations = [
        '• Deploy model in production environment with real-time scoring',
        '• Implement A/B testing for model-driven interventions',
        '• Collect additional features (device type, time of day, user history)',
        '• Develop customer segmentation using unsupervised learning',
        '• Create automated retraining pipeline with new data',
        '• Build comprehensive monitoring and alerting system'
    ]

    for rec in recommendations:
        p = tf.add_paragraph()
        p.text = rec
        p.level = 0

    # Save presentation
    prs.save('online_shopping_ml_presentation.pptx')
    print("Presentation saved as 'online_shopping_ml_presentation.pptx'")

if __name__ == "__main__":
    create_presentation()